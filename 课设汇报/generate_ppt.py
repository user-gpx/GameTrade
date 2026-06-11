#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
GameTrade 课程设计汇报 PPT 生成脚本
生成 18 页课堂答辩风格演示文稿，约 15-20 分钟汇报。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap


# ============================================================================
# 全局常量
# ============================================================================

# 幻灯片尺寸 (16:9 widescreen)
SLIDE_WIDTH = Cm(33.867)   # 13.333 inches
SLIDE_HEIGHT = Cm(19.05)   # 7.5 inches

# 配色方案
COLOR_PRIMARY = RGBColor(0x2c, 0x3e, 0x50)      # 深蓝灰 #2c3e50
COLOR_SECONDARY = RGBColor(0x34, 0x98, 0xdb)    # 蓝 #3498db
COLOR_ACCENT = RGBColor(0xe7, 0x4c, 0x3c)       # 红 #e74c3c
COLOR_BG = RGBColor(0xec, 0xf0, 0xf1)           # 浅灰 #ecf0f1
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
COLOR_LIGHT_GRAY = RGBColor(0xbd, 0xbd, 0xbd)
COLOR_GREEN = RGBColor(0x27, 0xae, 0x60)
COLOR_ORANGE = RGBColor(0xe6, 0x7e, 0x22)
COLOR_CHECK_GREEN = RGBColor(0x27, 0xae, 0x60)

# 正文字号
FONT_TITLE = Pt(32)
FONT_SUBTITLE = Pt(20)
FONT_SECTION_TITLE = Pt(28)
FONT_BODY = Pt(16)
FONT_SMALL = Pt(13)
FONT_CAPTION = Pt(11)
FONT_PAGE_NUM = Pt(10)

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, 'GameTrade课程设计汇报.pptx')


# ============================================================================
# 辅助函数
# ============================================================================

def set_slide_bg(slide, color=COLOR_BG):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_page_number(slide, page_num):
    """在右下角添加页码"""
    left = SLIDE_WIDTH - Cm(2.0)
    top = SLIDE_HEIGHT - Cm(1.0)
    width = Cm(1.5)
    height = Cm(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = FONT_PAGE_NUM
    p.font.color.rgb = COLOR_DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT


def add_bottom_line(slide):
    """在页面底部添加装饰线"""
    left = Cm(1.0)
    top = SLIDE_HEIGHT - Cm(1.3)
    width = SLIDE_WIDTH - Cm(2.0)
    height = Cm(0.03)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_SECONDARY
    shape.line.fill.background()


def add_slide_title(slide, title_text):
    """在页面顶部添加统一的标题栏"""
    # 标题背景条
    left = Cm(0)
    top = Cm(0)
    width = SLIDE_WIDTH
    height = Cm(2.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()

    # 标题下方细线
    left2 = Cm(1.5)
    top2 = Cm(2.8)
    width2 = SLIDE_WIDTH - Cm(3.0)
    height2 = Cm(0.06)
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left2, top2, width2, height2)
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLOR_SECONDARY
    shape2.line.fill.background()

    # 标题文字
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.5), width2, Cm(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = FONT_SECTION_TITLE
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT


def add_textbox(slide, left, top, width, height, text, font_size=FONT_BODY,
                color=COLOR_BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    """添加文本框的便捷函数"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=FONT_BODY,
                          color=COLOR_BLACK, bold_first=False, alignment=PP_ALIGN.LEFT,
                          line_spacing=Pt(24)):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.alignment = alignment
        p.space_after = line_spacing
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color=COLOR_WHITE,
                     border_color=None, border_width=Pt(0)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_placeholder_box(slide, left, top, width, height, text,
                        dashed=False, bg_color=None):
    """添加占位符矩形框（虚线边框）"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if bg_color is None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color

    if dashed:
        shape.line.color.rgb = COLOR_LIGHT_GRAY
        shape.line.width = Pt(1.5)
        shape.line.dash_style = 3  # dash style
    else:
        shape.line.color.rgb = COLOR_LIGHT_GRAY
        shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = FONT_SMALL
    p.font.color.rgb = COLOR_DARK_GRAY
    p.font.name = 'Microsoft YaHei'
    return shape


def add_card(slide, left, top, width, height, title, content_lines,
             title_color=COLOR_PRIMARY, card_bg=COLOR_WHITE):
    """添加卡片式布局（标题 + 内容）"""
    shape = add_rounded_rect(slide, left, top, width, height,
                             fill_color=card_bg, border_color=RGBColor(0xdd, 0xdd, 0xdd),
                             border_width=Pt(0.5))
    # 卡片标题线
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Cm(0.5), top + Cm(1.6),
                                    width - Cm(1.0), Cm(0.04))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = title_color
    shape2.line.fill.background()

    # 标题
    add_textbox(slide, left + Cm(0.7), top + Cm(0.3), width - Cm(1.4), Cm(1.2),
                title, font_size=FONT_SUBTITLE, color=title_color, bold=True)

    # 内容
    if content_lines:
        add_multiline_textbox(slide, left + Cm(0.7), top + Cm(2.0), width - Cm(1.4),
                              height - Cm(2.5), content_lines, font_size=FONT_SMALL,
                              color=COLOR_DARK_GRAY, line_spacing=Pt(18))


def add_decorative_bar(slide, left, top, width, height, color=COLOR_SECONDARY):
    """添加装饰性色条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ============================================================================
# 幻灯片创建函数
# ============================================================================

def create_cover_slide(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, COLOR_PRIMARY)

    # 顶部装饰线
    add_decorative_bar(slide, Cm(0), Cm(0), SLIDE_WIDTH, Cm(0.15), COLOR_SECONDARY)

    # 大标题
    add_textbox(slide, Cm(2.5), Cm(4.0), SLIDE_WIDTH - Cm(5.0), Cm(3.0),
                '游戏道具交易平台', font_size=Pt(42), color=COLOR_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 装饰线
    add_decorative_bar(slide, Cm(10.0), Cm(7.2), Cm(13.867), Cm(0.06), COLOR_SECONDARY)

    # 副标题
    add_textbox(slide, Cm(2.5), Cm(7.6), SLIDE_WIDTH - Cm(5.0), Cm(2.0),
                '课程设计汇报', font_size=Pt(28), color=COLOR_SECONDARY, bold=False,
                alignment=PP_ALIGN.CENTER)

    # 英文副标题
    add_textbox(slide, Cm(2.5), Cm(9.5), SLIDE_WIDTH - Cm(5.0), Cm(1.5),
                'GameTrade - Game Item Trading Platform', font_size=Pt(18),
                color=RGBColor(0x99, 0xbb, 0xcc), alignment=PP_ALIGN.CENTER)

    # 底部装饰线
    add_decorative_bar(slide, Cm(10.0), Cm(11.2), Cm(13.867), Cm(0.04), COLOR_SECONDARY)

    # 成员信息
    add_multiline_textbox(slide, Cm(8.0), Cm(12.0), Cm(18.0), Cm(5.0), [
        '成员A (曾雨111):  用户认证 + 道具展示',
        '成员B:  交易模块 + 支付系统',
        '成员C:  订单管理 + 统计功能 + 前端整合',
    ], font_size=FONT_SMALL, color=RGBColor(0xbb, 0xcc, 0xdd), line_spacing=Pt(16))

    # 日期
    add_textbox(slide, Cm(8.0), Cm(16.5), Cm(18.0), Cm(1.0),
                '2026年6月', font_size=FONT_SMALL, color=RGBColor(0x88, 0x99, 0xaa),
                alignment=PP_ALIGN.CENTER)

    # 右下角装饰三角
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, SLIDE_WIDTH - Cm(4.0),
                                   SLIDE_HEIGHT - Cm(4.0), Cm(4.0), Cm(4.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x3d, 0x54, 0x6a)
    shape.line.fill.background()
    # rotate
    shape.rotation = 180.0

    add_page_number(slide, 1)


def create_toc_slide(prs):
    """第2页：目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)

    # 左侧装饰条
    add_decorative_bar(slide, Cm(0), Cm(0), Cm(0.3), SLIDE_HEIGHT, COLOR_PRIMARY)

    add_textbox(slide, Cm(1.5), Cm(0.8), Cm(20.0), Cm(2.0),
                '目  录', font_size=Pt(36), color=COLOR_PRIMARY, bold=True)

    add_decorative_bar(slide, Cm(1.5), Cm(3.0), Cm(5.0), Cm(0.06), COLOR_SECONDARY)

    # 两列布局
    toc_items = [
        ('01', '项目背景', '游戏道具交易市场现状与痛点分析'),
        ('02', '需求分析', '功能需求与非功能需求规格说明'),
        ('03', '系统建模', '用例图、活动图、类图、时序图、状态图'),
        ('04', '概要设计', '体系结构设计与数据库设计'),
        ('05', '功能演示', '各模块功能实现与截图展示'),
        ('06', '测试与质量保证', '测试策略、用例与运行结果'),
        ('07', '总结与展望', '已完成工作总结与后续计划'),
    ]

    left_col_x = Cm(1.5)
    right_col_x = Cm(17.5)
    col_width = Cm(15.0)
    start_y = Cm(3.8)
    row_gap = Cm(2.2)

    for i, (num, title, desc) in enumerate(toc_items):
        if i <= 3:
            x = left_col_x
            y = start_y + i * row_gap
        else:
            x = right_col_x
            y = start_y + (i - 4) * row_gap

        # 编号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Cm(0.1), Cm(1.2), Cm(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_SECONDARY if i < 6 else COLOR_PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 标题
        add_textbox(slide, x + Cm(1.6), y, col_width - Cm(1.6), Cm(1.0),
                    title, font_size=Pt(18), color=COLOR_PRIMARY, bold=True)
        # 描述
        add_textbox(slide, x + Cm(1.6), y + Cm(1.0), col_width - Cm(1.6), Cm(1.0),
                    desc, font_size=FONT_SMALL, color=COLOR_DARK_GRAY)

    add_bottom_line(slide)
    add_page_number(slide, 2)


def create_background_slide(prs):
    """第3页：项目背景与业务需求"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '项目背景与业务需求')

    # 左侧卡片 - 市场痛点
    add_card(slide, Cm(1.0), Cm(3.5), Cm(15.5), Cm(14.0),
             '市场痛点',
             ['交易不安全：缺乏第三方担保，诈骗频发',
              '信息不对称：买卖双方难以及时匹配需求',
              '流程繁琐：传统交易需多方沟通，效率低下',
              '缺乏信任：卖家信誉无法验证，买家顾虑多',
              '支付风险：直接转账无保障，退款困难',
              '售后缺失：道具交付后出现问题无法追溯'],
             title_color=COLOR_ACCENT)

    # 右侧卡片 - 解决方案
    add_card(slide, Cm(17.5), Cm(3.5), Cm(15.5), Cm(14.0),
             '本平台解决方案',
             ['一站式交易：发布、浏览、购买全流程闭环',
              '安全可靠：用户认证 + 平台担保交易',
              '简洁高效：直观的UI设计，三步完成购买',
              '信誉体系：卖家信息透明，交易记录可查',
              '模拟支付：安全的余额系统，资金可追溯',
              '订单追踪：从下单到收货全流程状态追踪'],
             title_color=COLOR_SECONDARY)

    # 底部标语
    add_textbox(slide, Cm(1.0), Cm(17.8), SLIDE_WIDTH - Cm(2.0), Cm(1.0),
                '打造安全、便捷、高效的游戏道具交易体验',
                font_size=Pt(16), color=COLOR_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_bottom_line(slide)
    add_page_number(slide, 3)


def create_usecase_slide(prs):
    """第4页：需求分析 — 用例图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '需求分析 — 系统用例图')

    # 用例图占位符
    add_placeholder_box(slide, Cm(1.5), Cm(3.5), Cm(30.867), Cm(11.0),
                        '【用例图】\n见 diagrams/use_case.drawio\n请导出PNG后替换此占位符',
                        dashed=True)

    # 底部说明
    add_multiline_textbox(slide, Cm(1.5), Cm(15.5), Cm(30.0), Cm(3.0), [
        '三类用户角色：',
        '  买家（普通用户）：注册登录、浏览道具、搜索筛选、收藏道具、下单购买、确认收货、查看订单',
        '  卖家：发布道具、管理道具（编辑/下架）、发货操作、查看销售订单',
        '  管理员：用户管理、道具审核、数据统计',
    ], font_size=FONT_SMALL, color=COLOR_DARK_GRAY, line_spacing=Pt(16))

    add_bottom_line(slide)
    add_page_number(slide, 4)


def create_requirements_slide(prs):
    """第5页：需求规格说明"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '需求规格说明')

    # 左栏 - 功能需求
    shape_left = add_rounded_rect(slide, Cm(1.0), Cm(3.5), Cm(16.0), Cm(14.0),
                                  fill_color=COLOR_WHITE, border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                  border_width=Pt(0.5))
    add_textbox(slide, Cm(2.0), Cm(3.8), Cm(14.0), Cm(1.2),
                '功能需求', font_size=FONT_SUBTITLE, color=COLOR_SECONDARY, bold=True)
    add_decorative_bar(slide, Cm(2.0), Cm(5.0), Cm(4.0), Cm(0.04), COLOR_SECONDARY)

    func_items = [
        '1. 用户注册与登录（含个人资料编辑）',
        '2. 道具发布、编辑与下架管理',
        '3. 道具搜索、筛选与分页浏览',
        '4. 道具收藏与取消收藏（AJAX）',
        '5. 在线购买与下单',
        '6. 订单管理与状态追踪',
        '7. 余额充值（模拟支付）',
        '8. 收货确认与交易完成',
    ]
    add_multiline_textbox(slide, Cm(2.0), Cm(5.4), Cm(14.0), Cm(11.0), func_items,
                          font_size=FONT_SMALL, color=COLOR_BLACK, line_spacing=Pt(18))

    # 右栏 - 非功能需求
    shape_right = add_rounded_rect(slide, Cm(17.5), Cm(3.5), Cm(15.5), Cm(14.0),
                                   fill_color=COLOR_WHITE, border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                   border_width=Pt(0.5))
    add_textbox(slide, Cm(18.5), Cm(3.8), Cm(13.5), Cm(1.2),
                '非功能需求', font_size=FONT_SUBTITLE, color=COLOR_ACCENT, bold=True)
    add_decorative_bar(slide, Cm(18.5), Cm(5.0), Cm(4.0), Cm(0.04), COLOR_ACCENT)

    nonfunc_items = [
        '安全性 (Security)',
        '  - 用户认证与权限控制',
        '  - CSRF 防护，防止跨站攻击',
        '',
        '可用性 (Usability)',
        '  - Bootstrap 5 响应式设计',
        '  - 直观的导航与操作流程',
        '',
        '可靠性 (Reliability)',
        '  - 数据库事务保证数据一致性',
        '  - 订单状态机防止非法流转',
        '',
        '可维护性 (Maintainability)',
        '  - Django MTV 分层架构',
        '  - 模块化 App 设计，高内聚低耦合',
    ]
    add_multiline_textbox(slide, Cm(18.5), Cm(5.4), Cm(13.5), Cm(11.0), nonfunc_items,
                          font_size=FONT_SMALL, color=COLOR_BLACK, line_spacing=Pt(12))

    add_bottom_line(slide)
    add_page_number(slide, 5)


def create_activity_slide(prs):
    """第6页：系统建模 — 活动图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '系统建模 — 购买道具活动图')

    add_placeholder_box(slide, Cm(1.5), Cm(3.5), Cm(30.867), Cm(11.0),
                        '【活动图】\n见 diagrams/activity_buy_item.drawio\n请导出PNG后替换此占位符',
                        dashed=True)

    add_multiline_textbox(slide, Cm(1.5), Cm(15.5), Cm(30.0), Cm(3.0), [
        '关键流程节点：',
        '  用户浏览道具列表 → 查看详情 → 点击购买 → 创建订单 → 进入支付页面',
        '  → 选择支付方式（余额支付）→ 输入密码确认 → 支付成功 → 卖家发货 → 买家确认收货 → 交易完成',
    ], font_size=FONT_SMALL, color=COLOR_DARK_GRAY, line_spacing=Pt(16))

    add_bottom_line(slide)
    add_page_number(slide, 6)


def create_class_diagram_slide(prs):
    """第7页：系统建模 — 类图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '系统建模 — 核心领域模型类图')

    add_placeholder_box(slide, Cm(1.5), Cm(3.5), Cm(30.867), Cm(11.0),
                        '【类图】\n见 diagrams/class_diagram.drawio\n请导出PNG后替换此占位符',
                        dashed=True)

    add_multiline_textbox(slide, Cm(1.5), Cm(15.5), Cm(30.0), Cm(3.0), [
        '6个核心实体及其关系：',
        '  User ──1:1── UserProfile（用户扩展信息：头像、手机号、余额）',
        '  User ──1:N── Item（卖家发布道具）',
        '  Category ──1:N── Item（道具分类）',
        '  User ──M:N── Item (through Favorite)（收藏关系）',
        '  User (买家) ──1:N── Order ──N:1── User (卖家)（订单关系）',
        '  Order ──1:1── Payment ──1:N── TransactionLog（支付流水）',
    ], font_size=FONT_SMALL, color=COLOR_DARK_GRAY, line_spacing=Pt(16))

    add_bottom_line(slide)
    add_page_number(slide, 7)


def create_sequence_slide(prs):
    """第8页：系统建模 — 时序图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '系统建模 — 购买交易时序图')

    add_placeholder_box(slide, Cm(1.5), Cm(3.5), Cm(30.867), Cm(11.0),
                        '【时序图】\n见 diagrams/sequence_buy_trade.drawio\n请导出PNG后替换此占位符',
                        dashed=True)

    add_multiline_textbox(slide, Cm(1.5), Cm(15.5), Cm(30.0), Cm(3.0), [
        '关键交互步骤：',
        '  买家→浏览器→View→Model→模板：浏览道具列表 → 点击购买 → POST下单请求',
        '  → 创建Order记录(状态: pending) → 跳转支付页 → 确认支付 → 扣减余额',
        '  → 创建Payment+TransactionLog → 更新Order状态(paid) → 通知卖家 → 卖家发货',
        '  → 更新Order状态(shipped) → 买家确认收货 → 更新Order(completed) → 增加卖家余额',
    ], font_size=FONT_SMALL, color=COLOR_DARK_GRAY, line_spacing=Pt(16))

    add_bottom_line(slide)
    add_page_number(slide, 8)


def create_state_slide(prs):
    """第9页：系统建模 — 状态图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '系统建模 — 订单状态机')

    add_placeholder_box(slide, Cm(1.5), Cm(3.5), Cm(30.867), Cm(8.0),
                        '【状态图】\n见 diagrams/state_order.drawio\n请导出PNG后替换此占位符',
                        dashed=True)

    # 5种订单状态说明
    states = [
        ('pending', '待支付', '买家已创建订单，等待完成支付', COLOR_DARK_GRAY),
        ('paid', '已支付', '买家完成支付，等待卖家发货', COLOR_SECONDARY),
        ('shipped', '已发货', '卖家已发货，等待买家确认收货', COLOR_ORANGE),
        ('completed', '已完成', '买家确认收货，交易成功完成', COLOR_GREEN),
        ('cancelled', '已取消', '订单被取消（买家主动取消或超时）', COLOR_ACCENT),
    ]

    state_y = Cm(12.5)
    for i, (code, name, desc, clr) in enumerate(states):
        x_pos = Cm(1.5) + i * Cm(6.5)
        # 状态圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Cm(1.5), state_y, Cm(0.6), Cm(0.6))
        dot.fill.solid()
        dot.fill.fore_color.rgb = clr
        dot.line.fill.background()

        add_textbox(slide, x_pos, state_y + Cm(1.0), Cm(6.0), Cm(1.5),
                    f'{name} ({code})', font_size=FONT_SMALL, color=clr, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x_pos, state_y + Cm(2.0), Cm(6.0), Cm(1.0),
                    desc, font_size=FONT_CAPTION, color=COLOR_DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_bottom_line(slide)
    add_page_number(slide, 9)


def create_architecture_slide(prs):
    """第10页：概要设计 — 体系结构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '概要设计 — Django MTV 体系结构')

    # ---- 左侧：三层架构图 ----
    layer_data = [
        ('Template 层', 'Bootstrap 5 + jQuery + Font Awesome\nDjango 模板继承 (base.html)', COLOR_SECONDARY),
        ('View 层', '函数视图 + @login_required 装饰器\nAJAX 异步交互 + 消息框架', RGBColor(0x29, 0x80, 0xb9)),
        ('Model 层', '6 个 App 模块\nusers / items / trading / orders / payments / stats', COLOR_PRIMARY),
        ('数据库', 'SQLite（开发环境）/ PostgreSQL（生产环境）', RGBColor(0x1a, 0x25, 0x2f)),
    ]

    layer_y = Cm(3.5)
    layer_h = Cm(3.0)
    layer_w = Cm(16.0)
    for i, (name, desc, clr) in enumerate(layer_data):
        y = layer_y + i * (layer_h + Cm(0.8))
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), y, layer_w, layer_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = clr
        rect.line.fill.background()

        add_textbox(slide, Cm(2.5), y + Cm(0.3), layer_w - Cm(2.0), Cm(1.0),
                    name, font_size=FONT_SUBTITLE, color=COLOR_WHITE, bold=True)
        add_textbox(slide, Cm(2.5), y + Cm(1.3), layer_w - Cm(2.0), Cm(1.5),
                    desc, font_size=FONT_SMALL, color=RGBColor(0xee, 0xee, 0xee))

    # 添加层间箭头
    for i in range(3):
        y = layer_y + (i + 1) * (layer_h + Cm(0.8)) - Cm(0.7)
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Cm(9.0), y, Cm(1.0), Cm(0.6))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_SECONDARY
        arrow.line.fill.background()

    # ---- 右侧：模块列表 ----
    add_textbox(slide, Cm(19.0), Cm(3.5), Cm(13.0), Cm(1.5),
                '6 个 App 模块', font_size=FONT_SUBTITLE, color=COLOR_PRIMARY, bold=True)

    modules = [
        ('users', '用户认证', '注册/登录/登出/个人资料'),
        ('items', '道具管理', '发布/浏览/搜索/收藏'),
        ('trading', '交易执行', '购买流程/价格逻辑'),
        ('orders', '订单管理', '状态机/发货/收货'),
        ('payments', '支付处理', '模拟支付/流水记录'),
        ('stats', '统计分析', '月报/热门道具/图表'),
    ]

    mod_y = Cm(5.2)
    for i, (code, name, desc) in enumerate(modules):
        y = mod_y + i * Cm(2.2)
        rect = add_rounded_rect(slide, Cm(19.0), y, Cm(13.5), Cm(1.9),
                                fill_color=COLOR_WHITE,
                                border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                border_width=Pt(0.5))
        add_textbox(slide, Cm(19.5), y + Cm(0.1), Cm(4.0), Cm(0.8),
                    code, font_size=Pt(12), color=COLOR_SECONDARY, bold=True)
        add_textbox(slide, Cm(19.5), y + Cm(0.9), Cm(12.5), Cm(0.8),
                    f'{name} — {desc}', font_size=FONT_SMALL, color=COLOR_DARK_GRAY)

    add_bottom_line(slide)
    add_page_number(slide, 10)


def create_database_slide(prs):
    """第11页：数据库设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '数据库设计')

    # 表格数据
    db_tables = [
        ('User', 'id, username, password, email', 'Django内置用户模型', '—'),
        ('UserProfile', 'id, user(1:1), avatar, phone, bio, balance', 'User扩展信息', 'User (OneToOne)'),
        ('Category', 'id, name, icon', '道具分类', '—'),
        ('Item', 'id, name, game, price, description, image,\n  status, views_count, seller(FK), category(FK)', '核心道具实体', 'User (FK), Category (FK)'),
        ('Order', 'id, order_no, buyer(FK), seller(FK), item(FK),\n  amount, status, created_at', '交易订单', 'User(buyer/seller), Item'),
        ('Payment', 'id, order(1:1), amount, method, status,\n  paid_at', '支付记录', 'Order (OneToOne)'),
        ('TransactionLog', 'id, user(FK), order(FK), amount, type,\n  description, created_at', '交易流水', 'User, Order'),
        ('Favorite', 'id, user(FK), item(FK), created_at', '收藏关系', 'User, Item (unique_together)'),
    ]

    # 表头
    col_w = [Cm(4.0), Cm(13.5), Cm(5.0), Cm(8.5)]
    col_x = [Cm(1.0)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    header_y = Cm(3.5)
    header_h = Cm(1.2)

    # 表头背景
    for j, (w, header_text) in enumerate(zip(col_w, ['表名', '核心字段', '用途说明', '关联关系'])):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[j], header_y, w, header_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_PRIMARY
        rect.line.color.rgb = RGBColor(0x1a, 0x25, 0x2f)
        rect.line.width = Pt(0.5)
        tf = rect.text_frame
        tf.paragraphs[0].text = header_text
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 表体
    row_h = Cm(1.6)
    for i, (tbl_name, fields, purpose, relations) in enumerate(db_tables):
        row_y = header_y + header_h + i * row_h

        # 交替行颜色
        bg = COLOR_WHITE if i % 2 == 0 else RGBColor(0xf5, 0xf7, 0xfa)
        data = [tbl_name, fields, purpose, relations]
        for j, (w, text) in enumerate(zip(col_w, data)):
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x[j], row_y, w, row_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = bg
            rect.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)
            rect.line.width = Pt(0.5)
            tf = rect.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = text
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = COLOR_BLACK
            tf.paragraphs[0].font.name = 'Microsoft YaHei'
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT

    # 底部备注
    add_textbox(slide, Cm(1.0), Cm(17.5), Cm(30.0), Cm(1.0),
                'OneToOne = 一对一  |  FK (ForeignKey) = 一对多  |  unique_together = 多对多唯一约束',
                font_size=FONT_CAPTION, color=COLOR_DARK_GRAY)

    add_bottom_line(slide)
    add_page_number(slide, 11)


def create_demo_auth_items_slide(prs):
    """第12页：功能演示① — 用户认证与道具展示（成员A）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '功能演示 — 用户认证与道具展示【已完成】')

    # 已完成功能清单
    check_items = [
        '用户注册 / 登录 / 登出',
        '个人资料编辑（头像上传）',
        '道具发布 / 编辑 / 下架（软删除）',
        '道具列表页（分页、排序、搜索/筛选）',
        '道具详情页（含浏览量统计）',
        '收藏 / 取消收藏（AJAX 异步交互）',
        '我的发布管理',
    ]

    check_y = Cm(3.5)
    for i, text in enumerate(check_items):
        y = check_y + i * Cm(1.5)
        # 勾号
        add_textbox(slide, Cm(1.5), y, Cm(1.0), Cm(1.0),
                    '✅', font_size=Pt(18), color=COLOR_GREEN)
        # 文字
        add_textbox(slide, Cm(3.0), y, Cm(18.0), Cm(1.0),
                    text, font_size=FONT_BODY, color=COLOR_BLACK)

    # 截图占位区域
    add_placeholder_box(slide, Cm(16.0), Cm(3.5), Cm(16.5), Cm(13.5),
                        '【截图占位】\n\n\n注册页面\n道具列表页\n道具详情页\n个人中心',
                        dashed=True, bg_color=RGBColor(0xfa, 0xfa, 0xfa))

    # 底部标注
    add_textbox(slide, Cm(1.5), Cm(17.5), Cm(14.0), Cm(1.0),
                '成员A: 用户认证 + 道具展示模块  ✅',
                font_size=FONT_SMALL, color=COLOR_PRIMARY, bold=True)

    add_bottom_line(slide)
    add_page_number(slide, 12)


def create_demo_trading_slide(prs):
    """第13页：功能演示② — 交易流程（成员B）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '功能演示 — 交易流程【已完成】')

    check_items = [
        '购买下单（buy_now 一键购买）',
        '创建订单 / 取消订单',
        '模拟支付 / 支付回调处理',
        '卖家发货操作',
        '买家确认收货',
        '交易流水记录（TransactionLog）',
        '买家 / 卖家订单列表',
    ]

    check_y = Cm(3.5)
    for i, text in enumerate(check_items):
        y = check_y + i * Cm(1.5)
        add_textbox(slide, Cm(1.5), y, Cm(1.0), Cm(1.0),
                    '✅', font_size=Pt(18), color=COLOR_GREEN)
        add_textbox(slide, Cm(3.0), y, Cm(18.0), Cm(1.0),
                    text, font_size=FONT_BODY, color=COLOR_BLACK)

    add_placeholder_box(slide, Cm(16.0), Cm(3.5), Cm(16.5), Cm(13.5),
                        '【截图占位】\n\n\n购买页面\n支付确认页\n发货页面\n订单列表页',
                        dashed=True, bg_color=RGBColor(0xfa, 0xfa, 0xfa))

    add_textbox(slide, Cm(1.5), Cm(17.5), Cm(14.0), Cm(1.0),
                '成员B: 交易模块 + 支付系统  ✅',
                font_size=FONT_SMALL, color=COLOR_PRIMARY, bold=True)

    add_bottom_line(slide)
    add_page_number(slide, 13)


def create_demo_recharge_slide(prs):
    """第14页：功能演示③ — 余额充值（成员B）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '功能演示 — 余额充值【简化版已完成】')

    # 说明
    add_multiline_textbox(slide, Cm(1.5), Cm(3.5), Cm(16.0), Cm(3.0), [
        '当前实现：模拟充值，直接增加用户余额数字，未接入真实支付接口。',
        '完整的 TransactionLog 流水记录，确保每笔资金变动可追溯。',
    ], font_size=FONT_BODY, color=COLOR_DARK_GRAY, line_spacing=Pt(22))

    # 充值流程
    add_textbox(slide, Cm(1.5), Cm(6.5), Cm(15.0), Cm(1.0),
                '充值流程', font_size=FONT_SUBTITLE, color=COLOR_PRIMARY, bold=True)

    steps = [
        ('Step 1', '用户点击充值', '访问充值页面'),
        ('Step 2', '输入充值金额', '填写金额并确认'),
        ('Step 3', '系统增加余额', '直接增加余额数字'),
        ('Step 4', '记录流水', '写入 TransactionLog'),
    ]

    step_x_start = Cm(1.5)
    step_y = Cm(8.0)
    step_w = Cm(6.5)
    step_h = Cm(5.5)
    for i, (step, title, desc) in enumerate(steps):
        x = step_x_start + i * (step_w + Cm(1.0))
        rect = add_rounded_rect(slide, x, step_y, step_w, step_h,
                                fill_color=COLOR_WHITE,
                                border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                border_width=Pt(0.5))
        add_textbox(slide, x + Cm(0.5), step_y + Cm(0.3), step_w - Cm(1.0), Cm(1.0),
                    step, font_size=Pt(22), color=COLOR_SECONDARY, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.5), step_y + Cm(2.0), step_w - Cm(1.0), Cm(1.0),
                    title, font_size=FONT_BODY, color=COLOR_PRIMARY, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(0.5), step_y + Cm(3.2), step_w - Cm(1.0), Cm(1.0),
                    desc, font_size=FONT_SMALL, color=COLOR_DARK_GRAY, alignment=PP_ALIGN.CENTER)

        # 箭头 (不是最后一步)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + step_w + Cm(0.1), step_y + Cm(2.3),
                                           Cm(0.8), Cm(0.8))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_SECONDARY
            arrow.line.fill.background()

    # 截图占位
    add_placeholder_box(slide, Cm(1.5), Cm(14.2), Cm(30.0), Cm(2.5),
                        '【截图占位：充值页面】',
                        dashed=True)

    # 底部备注
    add_textbox(slide, Cm(1.5), Cm(17.5), Cm(30.0), Cm(1.0),
                'ℹ 后续可接入支付宝/微信支付沙箱环境，实现真实支付流程',
                font_size=FONT_SMALL, color=COLOR_ACCENT)

    add_page_number(slide, 14)


def create_demo_orders_slide(prs):
    """第15页：功能演示④ — 订单管理页面（成员C）【占位】"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '功能演示 — 订单管理页面【待合并】')

    # 大占位框 - 虚线边框 + 浅灰背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3.0), Cm(3.8), Cm(28.0), Cm(13.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 3  # dash

    add_multiline_textbox(slide, Cm(5.0), Cm(4.5), Cm(24.0), Cm(12.0), [
        '【占位：待合并成员C代码】',
        '',
        '预计内容包括：',
        '',
        '  •  买家订单列表页（按状态筛选：待支付/已支付/已发货/已完成/已取消）',
        '  •  卖家订单列表页（发货操作按钮）',
        '  •  订单详情页（完整订单信息展示）',
        '',
        '',
        '',
        '',
        '（虚线边框 + 浅灰背景标记为占位区域）',
    ], font_size=FONT_BODY, color=COLOR_DARK_GRAY, line_spacing=Pt(20),
       alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Cm(1.5), Cm(17.8), Cm(14.0), Cm(1.0),
                '成员C: 订单管理  \U0001F532',
                font_size=FONT_SMALL, color=COLOR_PRIMARY, bold=True)

    add_bottom_line(slide)
    add_page_number(slide, 15)


def create_demo_stats_slide(prs):
    """第16页：功能演示⑤ — 统计模块与前端美化（成员C）【占位】"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '功能演示 — 统计模块与前端整合【待合并】')

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3.0), Cm(3.8), Cm(28.0), Cm(13.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xf8, 0xf8)
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 3

    add_multiline_textbox(slide, Cm(5.0), Cm(4.5), Cm(24.0), Cm(12.0), [
        '【占位：待合并成员C代码】',
        '',
        '预计内容包括：',
        '',
        '  •  热门道具统计 / 月报生成',
        '  •  邮件发送功能（月报推送给用户）',
        '  •  前端整体 Bootstrap 美化',
        '  •  响应式布局优化（移动端适配）',
        '  •  全局导航栏与页面风格统一',
        '',
        '',
        '',
        '（虚线边框 + 浅灰背景标记为占位区域）',
    ], font_size=FONT_BODY, color=COLOR_DARK_GRAY, line_spacing=Pt(20),
       alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Cm(1.5), Cm(17.8), Cm(14.0), Cm(1.0),
                '成员C: 统计 + 前端整合  \U0001F532',
                font_size=FONT_SMALL, color=COLOR_PRIMARY, bold=True)

    add_bottom_line(slide)
    add_page_number(slide, 16)


def create_testing_slide(prs):
    """第17页：测试与质量保证"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '测试与质量保证')

    # 左栏 - 测试类型
    shape_left = add_rounded_rect(slide, Cm(1.0), Cm(3.5), Cm(16.5), Cm(14.0),
                                  fill_color=COLOR_WHITE,
                                  border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                  border_width=Pt(0.5))
    add_textbox(slide, Cm(2.0), Cm(3.8), Cm(14.0), Cm(1.0),
                '测试策略', font_size=FONT_SUBTITLE, color=COLOR_SECONDARY, bold=True)
    add_decorative_bar(slide, Cm(2.0), Cm(4.8), Cm(3.0), Cm(0.04), COLOR_SECONDARY)

    test_lines = [
        '单元测试 (Unit Test)',
        '  • 用户认证：注册/登录/登出',
        '  • 道具浏览：列表/详情/搜索/筛选',
        '  • 交易流程：购买/支付/发货/收货',
        '  • 收藏功能：添加/取消收藏',
        '',
        '黑盒测试 (Black-box Test)',
        '  • 功能测试用例设计',
        '  • 等价类划分法',
        '  • 边界值分析法',
        '',
        '白盒测试 (White-box Test)',
        '  • 语句覆盖（Statement Coverage）',
        '  • 分支覆盖（Branch Coverage）',
        '',
        '自动化测试',
        '  • Selenium WebDriver UI 交互测试',
        '  • Django TestCase 集成测试',
    ]
    add_multiline_textbox(slide, Cm(2.0), Cm(5.2), Cm(14.5), Cm(11.0),
                          test_lines, font_size=FONT_SMALL, color=COLOR_BLACK,
                          line_spacing=Pt(10))

    # 右栏 - 测试统计
    shape_right = add_rounded_rect(slide, Cm(18.5), Cm(3.5), Cm(14.5), Cm(14.0),
                                   fill_color=COLOR_WHITE,
                                   border_color=RGBColor(0xdd, 0xdd, 0xdd),
                                   border_width=Pt(0.5))
    add_textbox(slide, Cm(19.5), Cm(3.8), Cm(12.5), Cm(1.0),
                '测试统计', font_size=FONT_SUBTITLE, color=COLOR_PRIMARY, bold=True)
    add_decorative_bar(slide, Cm(19.5), Cm(4.8), Cm(3.0), Cm(0.04), COLOR_PRIMARY)

    add_placeholder_box(slide, Cm(20.0), Cm(5.5), Cm(11.5), Cm(10.0),
                        '【截图占位】\n\n测试运行结果\n\n\npython manage.py test\n\n测试覆盖率报告',
                        dashed=True, bg_color=RGBColor(0xf5, 0xf5, 0xf5))

    add_bottom_line(slide)
    add_page_number(slide, 17)


def create_summary_slide(prs):
    """第18页：总结与展望"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_BG)
    add_slide_title(slide, '总结与展望')

    # 上半：已完成工作
    add_textbox(slide, Cm(1.5), Cm(3.3), Cm(15.0), Cm(2.0),
                '已完成工作', font_size=FONT_SUBTITLE, color=COLOR_GREEN, bold=True)

    # 成员A 卡片
    add_card(slide, Cm(1.5), Cm(5.0), Cm(15.5), Cm(4.5),
             '成员A', [
                 '用户认证模块（注册/登录/登出/个人资料编辑）',
                 '道具展示模块（发布/列表/详情/搜索/收藏）',
                 '状态：✅ 已完成',
             ], title_color=COLOR_SECONDARY)

    # 成员B 卡片
    add_card(slide, Cm(17.5), Cm(5.0), Cm(15.5), Cm(4.5),
             '成员B', [
                 '交易核心模块（购买/支付/发货/收货流程）',
                 '支付系统（模拟支付/余额充值/交易流水）',
                 '状态：✅ 已完成',
             ], title_color=COLOR_SECONDARY)

    # 下半：待完成工作
    add_textbox(slide, Cm(1.5), Cm(10.5), Cm(15.0), Cm(2.0),
                '待完成工作', font_size=FONT_SUBTITLE, color=COLOR_ACCENT, bold=True)

    # 成员C 卡片
    add_card(slide, Cm(1.5), Cm(12.0), Cm(15.5), Cm(4.5),
             '成员C', [
                 '订单管理页面（买家/卖家订单列表与详情）',
                 '统计模块（热门道具统计/月报生成/邮件发送）',
                 '前端整合（Bootstrap美化/响应式布局/全局导航）',
                 '状态：\U0001F532 待合并',
             ], title_color=COLOR_ORANGE)

    # 支付升级 卡片
    add_card(slide, Cm(17.5), Cm(12.0), Cm(15.5), Cm(4.5),
             '后续计划', [
                 '支付模块升级（对接支付宝/微信支付沙箱）',
                 '增加更多游戏类型支持',
                 '完善举报与纠纷处理机制',
                 '状态：\U0001F532 后续版本',
             ], title_color=COLOR_ACCENT)

    # 致谢
    add_textbox(slide, Cm(1.5), Cm(17.5), SLIDE_WIDTH - Cm(3.0), Cm(1.0),
                '感谢各位老师聆听！欢迎提问与指导。',
                font_size=FONT_SUBTITLE, color=COLOR_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_bottom_line(slide)
    add_page_number(slide, 18)


# ============================================================================
# 主函数
# ============================================================================

def main():
    """生成完整 PPT"""
    # 确保输出目录存在
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 按顺序创建所有幻灯片
    print('正在生成第 1 页：封面...')
    create_cover_slide(prs)

    print('正在生成第 2 页：目录...')
    create_toc_slide(prs)

    print('正在生成第 3 页：项目背景与业务需求...')
    create_background_slide(prs)

    print('正在生成第 4 页：需求分析 — 用例图...')
    create_usecase_slide(prs)

    print('正在生成第 5 页：需求规格说明...')
    create_requirements_slide(prs)

    print('正在生成第 6 页：系统建模 — 活动图...')
    create_activity_slide(prs)

    print('正在生成第 7 页：系统建模 — 类图...')
    create_class_diagram_slide(prs)

    print('正在生成第 8 页：系统建模 — 时序图...')
    create_sequence_slide(prs)

    print('正在生成第 9 页：系统建模 — 状态图...')
    create_state_slide(prs)

    print('正在生成第10页：概要设计 — 体系结构...')
    create_architecture_slide(prs)

    print('正在生成第11页：数据库设计...')
    create_database_slide(prs)

    print('正在生成第12页：功能演示① — 用户认证与道具展示...')
    create_demo_auth_items_slide(prs)

    print('正在生成第13页：功能演示② — 交易流程...')
    create_demo_trading_slide(prs)

    print('正在生成第14页：功能演示③ — 余额充值...')
    create_demo_recharge_slide(prs)

    print('正在生成第15页：功能演示④ — 订单管理页面...')
    create_demo_orders_slide(prs)

    print('正在生成第16页：功能演示⑤ — 统计模块与前端整合...')
    create_demo_stats_slide(prs)

    print('正在生成第17页：测试与质量保证...')
    create_testing_slide(prs)

    print('正在生成第18页：总结与展望...')
    create_summary_slide(prs)

    # 保存
    prs.save(OUTPUT_FILE)
    print(f'\nPPT 已成功生成！')
    print(f'文件路径：{OUTPUT_FILE}')
    print(f'共 {len(prs.slides)} 页幻灯片')


if __name__ == '__main__':
    main()
